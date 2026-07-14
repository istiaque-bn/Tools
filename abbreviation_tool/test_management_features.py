from django.contrib.auth import get_user_model
from django.core.files.uploadedfile import SimpleUploadedFile
from django.test import TestCase
from django.urls import reverse
from io import BytesIO
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches

from .models import AbbreviationEntry, AbbreviationProfile, Feedback
from .services.matching import POLICY_DEFINE_FIRST, candidates_for, find_matches
from .services.ooxml import CharacterLocation, TextContainer


class ManagementFeatureTests(TestCase):
    def setUp(self):
        self.admin = get_user_model().objects.create_user("admin2", password="test", is_staff=True)
        self.client.force_login(self.admin)

    def test_staff_can_add_edit_and_import(self):
        response = self.client.post(reverse("abbreviation_tool:manage_dictionary"), {"entry-abbreviation": "WHO", "entry-full_form": "World Health Organization", "save_entry": "1"})
        self.assertRedirects(response, reverse("abbreviation_tool:manage_dictionary"))
        entry = AbbreviationEntry.objects.get(abbreviation="WHO")
        self.assertTrue(entry.profiles.filter(name="General").exists())
        self.client.post(reverse("abbreviation_tool:edit_dictionary_entry", args=[entry.pk]), {"entry-abbreviation": "W.H.O.", "entry-full_form": "World Health Organization", "save_entry": "1"})
        entry.refresh_from_db()
        self.assertEqual(entry.abbreviation, "W.H.O.")
        upload = SimpleUploadedFile("entries.csv", b"abbreviation,full_form\nUN,United Nations\n", content_type="text/csv")
        self.client.post(reverse("abbreviation_tool:manage_dictionary"), {"import-file": upload, "import_entries": "1"})
        self.assertTrue(AbbreviationEntry.objects.filter(abbreviation="UN").exists())

    def test_feedback_is_stored(self):
        self.client.post(reverse("abbreviation_tool:feedback"), {"name": "Admin", "email": "admin@example.com", "message": "Please correct an entry."})
        self.assertEqual(Feedback.objects.count(), 1)

    def test_staff_can_delete_entry_and_sees_success(self):
        entry = AbbreviationEntry.objects.create(abbreviation="DEL", full_form="Delete Example")
        response = self.client.post(reverse("abbreviation_tool:delete_dictionary_entry", args=[entry.pk]), follow=True)
        self.assertFalse(AbbreviationEntry.objects.filter(pk=entry.pk).exists())
        self.assertContains(response, "was deleted successfully")

    def test_management_list_searches_abbreviation_and_full_form(self):
        AbbreviationEntry.objects.create(abbreviation="WHO", full_form="World Health Organization")
        AbbreviationEntry.objects.create(abbreviation="UN", full_form="United Nations")
        response = self.client.get(reverse("abbreviation_tool:manage_dictionary"), {"q": "WHO"})
        self.assertContains(response, "World Health Organization")
        self.assertEqual(list(response.context["search_results"].values_list("abbreviation", flat=True)), ["WHO"])
        response = self.client.get(reverse("abbreviation_tool:manage_dictionary"), {"q": "Nations"})
        self.assertContains(response, "United Nations")

    def test_add_form_keeps_action_when_submit_button_is_disabled(self):
        response = self.client.post(reverse("abbreviation_tool:manage_dictionary"), {
            "save_entry": "1",
            "entry-abbreviation": "NATO",
            "entry-full_form": "North Atlantic Treaty Organization",
        })
        self.assertRedirects(response, reverse("abbreviation_tool:manage_dictionary"))
        self.assertTrue(AbbreviationEntry.objects.filter(abbreviation="NATO").exists())

    def test_management_page_shows_total_entry_count(self):
        response = self.client.get(reverse("abbreviation_tool:manage_dictionary"))
        self.assertContains(response, "Total abbreviations")
        self.assertEqual(response.context["total_entries"], AbbreviationEntry.objects.count())

    def test_every_abbreviation_is_assigned_to_general(self):
        general = AbbreviationProfile.objects.get(name="General")
        self.assertEqual(general.entries.count(), AbbreviationEntry.objects.count())

    def test_staff_can_import_xlsx(self):
        workbook = Workbook()
        sheet = workbook.active
        sheet.append(["abbreviation", "full_form"])
        sheet.append(["EU", "European Union"])
        content = BytesIO()
        workbook.save(content)
        upload = SimpleUploadedFile("entries.xlsx", content.getvalue(), content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        response = self.client.post(reverse("abbreviation_tool:manage_dictionary"), {"import-file": upload, "import_entries": "1"})
        self.assertRedirects(response, reverse("abbreviation_tool:manage_dictionary"))
        self.assertTrue(AbbreviationEntry.objects.filter(abbreviation="EU", full_form="European Union").exists())

    def test_staff_can_export_every_abbreviation_to_xlsx(self):
        AbbreviationEntry.objects.create(abbreviation="WHO", full_form="World Health Organization")
        AbbreviationEntry.objects.create(abbreviation="UN", full_form="United Nations")

        response = self.client.get(reverse("abbreviation_tool:export_dictionary_xlsx"))

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response["Content-Type"], "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        self.assertIn(".xlsx", response["Content-Disposition"])
        workbook = load_workbook(BytesIO(b"".join(response.streaming_content)), read_only=True)
        rows = list(workbook.active.iter_rows(values_only=True))
        self.assertEqual(rows[0], ("abbreviation", "full_form"))
        self.assertIn(("WHO", "World Health Organization"), rows)
        self.assertIn(("UN", "United Nations"), rows)
        self.assertEqual(len(rows), AbbreviationEntry.objects.count() + 1)

    def test_multiple_full_forms_are_kept_as_separate_ambiguous_meanings(self):
        AbbreviationEntry.objects.create(abbreviation="CO", full_form="Commanding Officer")
        AbbreviationEntry.objects.create(abbreviation="CO", full_form="Company")
        response = self.client.get(reverse("abbreviation_tool:manage_dictionary"), {"q": "CO"})
        self.assertContains(response, "Commanding Officer")
        self.assertContains(response, "Company")
        self.assertContains(response, "2 meanings")

    def test_define_first_avoids_duplicate_parentheses(self):
        entry = AbbreviationEntry.objects.create(abbreviation="WHO", full_form="World Health Organization")
        text = "World Health Organization (WHO) works. World Health Organization reports."
        locations = [CharacterLocation(0, i, b"") for i in range(len(text))]
        matches = find_matches(TextContainer("text", "p", "paragraph", text, locations), candidates_for(None, "abbreviate"), "abbreviate", POLICY_DEFINE_FIRST)
        self.assertEqual(len(matches), 1)
        self.assertEqual(matches[0].proposed, "WHO")

    def test_text_converter_abbreviates_every_occurrence_without_definition(self):
        entry = AbbreviationEntry.objects.create(abbreviation="WHO", full_form="World Health Organization")
        general = AbbreviationProfile.objects.get(name="General")
        entry.profiles.add(general)
        response = self.client.post(reverse("abbreviation_tool:text_convert"), '{"operation":"abbreviate","text":"World Health Organization and World Health Organization"}', content_type="application/json")
        self.assertEqual(response.json()["result"], "WHO and WHO")

    def test_powerpoint_converter_preserves_presentation_and_defines_first_use(self):
        entry = AbbreviationEntry.objects.create(abbreviation="HQ", full_form="Headquarters")
        entry.profiles.add(AbbreviationProfile.objects.get(name="General"))
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        box.text_frame.paragraphs[0].text = "Headquarters and Headquarters"
        content = BytesIO()
        presentation.save(content)
        upload = SimpleUploadedFile("briefing.pptx", content.getvalue(), content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        response = self.client.post(reverse("abbreviation_tool:powerpoint_convert"), {"operation_type": "abbreviate", "presentation_file": upload})
        self.assertEqual(response.status_code, 200)
        converted = Presentation(BytesIO(b"".join(response.streaming_content)))
        self.assertEqual(converted.slides[0].shapes[0].text, "Headquarters (HQ) and HQ")

    def test_powerpoint_abbreviates_all_full_forms_even_when_abbreviation_is_ambiguous(self):
        general = AbbreviationProfile.objects.get(name="General")
        officer = AbbreviationEntry.objects.create(abbreviation="CO", full_form="Commanding Officer", is_ambiguous=True)
        company = AbbreviationEntry.objects.create(abbreviation="CO", full_form="Company", is_ambiguous=True)
        officer.profiles.add(general)
        company.profiles.add(general)
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        box.text_frame.paragraphs[0].text = "Commanding Officer met the Commanding Officer at the Company."
        content = BytesIO()
        presentation.save(content)
        upload = SimpleUploadedFile("ambiguous.pptx", content.getvalue(), content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")

        response = self.client.post(reverse("abbreviation_tool:powerpoint_convert"), {"operation_type": "abbreviate", "presentation_file": upload})

        self.assertEqual(response.status_code, 200)
        converted = Presentation(BytesIO(b"".join(response.streaming_content)))
        self.assertEqual(converted.slides[0].shapes[0].text, "Commanding Officer (CO) met the CO at the Company (CO).")

    def test_legacy_ppt_receives_safe_conversion_message(self):
        upload = SimpleUploadedFile("legacy.ppt", b"legacy", content_type="application/vnd.ms-powerpoint")
        response = self.client.post(reverse("abbreviation_tool:powerpoint_convert"), {"operation_type": "abbreviate", "presentation_file": upload})
        self.assertEqual(response.status_code, 400)
        self.assertContains(response, "save it as .pptx first", status_code=400)

    def test_longer_full_form_is_matched_before_nested_shorter_form(self):
        AbbreviationEntry.objects.create(abbreviation="BAF", full_form="Bangladesh Armed Force")
        AbbreviationEntry.objects.get_or_create(abbreviation="BD", full_form="Bangladesh")
        text = "Bangladesh Armed Force works for Bangladesh."
        locations = [CharacterLocation(0, i, b"") for i in range(len(text))]
        matches = find_matches(TextContainer("text", "p", "paragraph", text, locations), candidates_for(None, "abbreviate"), "abbreviate", POLICY_DEFINE_FIRST)
        self.assertEqual([match.original for match in matches], ["Bangladesh Armed Force", "Bangladesh"])
        self.assertEqual([match.proposed for match in matches], ["Bangladesh Armed Force (BAF)", "Bangladesh (BD)"])

    def test_managed_long_phrase_is_available_in_general_profile(self):
        self.client.post(reverse("abbreviation_tool:manage_dictionary"), {
            "save_entry": "1", "entry-abbreviation": "BDAF", "entry-full_form": "Bangladesh Armed Forces",
        })
        general = AbbreviationProfile.objects.get(name="General")
        text = "Bangladesh Armed Forces"
        locations = [CharacterLocation(0, i, b"") for i in range(len(text))]
        matches = find_matches(TextContainer("text", "p", "paragraph", text, locations), candidates_for(general, "abbreviate"), "abbreviate", POLICY_DEFINE_FIRST)
        self.assertEqual(len(matches), 1)
        self.assertEqual(matches[0].proposed, "Bangladesh Armed Forces (BDAF)")
