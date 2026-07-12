from io import BytesIO

from django.core.exceptions import ValidationError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from abbreviation_tool.models import DocumentProcessingSession

from .matching import POLICY_DEFINE_FIRST, candidates_for, find_matches
from .ooxml import CharacterLocation, TextContainer


def _paragraph_text(paragraph):
    return "".join(run.text for run in paragraph.runs)


def _replace_range(paragraph, start, end, replacement):
    runs = list(paragraph.runs)
    mapping = [(run_index, offset) for run_index, run in enumerate(runs) for offset in range(len(run.text))]
    if start < 0 or end > len(mapping) or start >= end:
        raise ValidationError("A PowerPoint text range could not be replaced safely.")
    first_run_index, first_offset = mapping[start]
    last_run_index, last_offset = mapping[end - 1]
    first_run = runs[first_run_index]
    last_run = runs[last_run_index]
    if first_run_index == last_run_index:
        first_run.text = first_run.text[:first_offset] + replacement + first_run.text[last_offset + 1:]
        return
    first_run.text = first_run.text[:first_offset] + replacement
    for run in runs[first_run_index + 1:last_run_index]:
        run.text = ""
    last_run.text = last_run.text[last_offset + 1:]


def _shape_paragraphs(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            yield from _shape_paragraphs(child)
    if getattr(shape, "has_text_frame", False):
        yield from shape.text_frame.paragraphs
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                yield from cell.text_frame.paragraphs


def process_powerpoint(upload, operation, profile):
    try:
        upload.seek(0)
        presentation = Presentation(upload)
    except Exception as exc:
        raise ValidationError("The PPTX file is invalid, encrypted, or damaged.") from exc
    candidates = candidates_for(profile, operation)
    seen = set()
    replacement_count = 0
    paragraph_index = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            for paragraph in _shape_paragraphs(shape):
                text = _paragraph_text(paragraph)
                if not text:
                    paragraph_index += 1
                    continue
                character_map = [CharacterLocation(run_index, offset, str(run_index).encode()) for run_index, run in enumerate(paragraph.runs) for offset in range(len(run.text))]
                container = TextContainer("pptx", f"pptx:p{paragraph_index}", "powerpoint_paragraph", text, character_map)
                matches = find_matches(container, candidates, operation, POLICY_DEFINE_FIRST, seen)
                matches = [match for match in matches if match.ambiguity == "unambiguous"]
                for match in sorted(matches, key=lambda item: item.start, reverse=True):
                    _replace_range(paragraph, match.start, match.end, match.proposed)
                    replacement_count += 1
                paragraph_index += 1
    output = BytesIO()
    presentation.save(output)
    output.seek(0)
    return output, replacement_count
